import sys

def extract_text(pptx_path):
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except ImportError:
        return "Error: python-pptx is not installed."
    except Exception as e:
        return f"Error extracting text from PPTX: {e}"

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx.py <path_to_pptx> <output_txt_path>")
        sys.exit(1)
    
    text = extract_text(sys.argv[1])
    with open(sys.argv[2], "w", encoding="utf-8") as f:
        f.write(text)
    print(f"Extraction saved to {sys.argv[2]}")
